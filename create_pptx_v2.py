from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- Design Configuration ---
THEME_COLOR_PRIMARY = RGBColor(0, 51, 102)    # Dark Navy
THEME_COLOR_SECONDARY = RGBColor(255, 102, 0) # Orange Accent
THEME_COLOR_BG = RGBColor(250, 250, 250)      # Light Gray Background
TEXT_COLOR_TITLE = RGBColor(0, 0, 0)
TEXT_COLOR_BODY = RGBColor(50, 50, 50)

def set_slide_background(slide):
    """Sets a subtle background design for the slide."""
    # Top bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.15)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = THEME_COLOR_PRIMARY
    shape.line.fill.background()

    # Bottom bar with accent
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(7.35), Inches(13.333), Inches(0.15)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = THEME_COLOR_PRIMARY
    shape.line.fill.background()
    
    # Accent mark
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(12), Inches(7.35), Inches(1), Inches(0.15)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = THEME_COLOR_SECONDARY
    shape.line.fill.background()

def create_title_slide(prs, title_text, subtitle_text):
    slide_layout = prs.slide_layouts[6] # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = THEME_COLOR_PRIMARY
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = title_text
    p.font.size = Pt(44)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11), Inches(2))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = subtitle_text
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.LEFT

def create_content_slide(prs, title, content_list, highlight_box=None):
    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    
    # Title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    tf = title_shape.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = THEME_COLOR_PRIMARY
    
    # Content Body
    body_shape = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(5))
    tf = body_shape.text_frame
    tf.word_wrap = True
    
    for item in content_list:
        p = tf.add_paragraph()
        
        # Check for sub-points (simple indentation logic)
        if item.startswith("   -"):
            p.text = item.strip()
            p.level = 1
            p.font.size = Pt(20)
        elif item.startswith("     *"):
            p.text = item.strip()
            p.level = 2
            p.font.size = Pt(18)
        else:
            p.text = item
            p.level = 0
            p.font.size = Pt(24)
            p.font.bold = True if ":" in item else False # Simple bolding rule
            p.space_before = Pt(12)
        
        p.font.color.rgb = TEXT_COLOR_BODY
        p.space_after = Pt(6)

    # Optional Highlight Box
    if highlight_box:
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8), Inches(2), Inches(4.5), Inches(3)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(240, 240, 240)
        box.line.color.rgb = THEME_COLOR_SECONDARY
        
        tf = box.text_frame
        tf.margin_top = Cm(0.5)
        tf.margin_left = Cm(0.5)
        p = tf.add_paragraph()
        p.text = highlight_box
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER

def create_comparison_table_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    
    # Title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    tf = title_shape.text_frame
    p = tf.add_paragraph()
    p.text = "Why This System? (핵심 경쟁력)"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = THEME_COLOR_PRIMARY
    
    # Table
    rows = 5
    cols = 3
    left = Inches(1)
    top = Inches(2)
    width = Inches(11.3)
    height = Inches(4)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Column widths
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(4.4)
    table.columns[2].width = Inches(4.4)
    
    # Headers
    headers = ["Feature", "일반 AI 이미지 생성 툴", "제안 시스템 (Fashion AI)"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = THEME_COLOR_PRIMARY if i != 2 else THEME_COLOR_SECONDARY
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Data
    data = [
        ["데이터 소스", "없음 (사용자 입력 의존)", "실시간 크롤링 (15+ 채널)"],
        ["분석 능력", "없음", "LLM 트렌드/감성 분석"],
        ["결과물", "단순 이미지 (JPG)", "이미지 + 도면(Blueprint) + 치수표"],
        ["현장 활용", "단순 아이디어 참고용", "샘플실 전달 가능한 작업지시서"]
    ]
    
    for r_idx, row_data in enumerate(data, 1):
        for c_idx, text in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.text = text
            cell.text_frame.paragraphs[0].font.size = Pt(16)
            if c_idx == 2:
                cell.text_frame.paragraphs[0].font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 240, 230) # Light orange bg

def create_presentation():
    prs = Presentation()
    # Widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Title Slide
    create_title_slide(
        prs, 
        "AI 기반 패션 트렌드 분석 및\n디자인 자동 생성 시스템 구축",
        "Data-Driven Fashion Design & Trend Forecasting System\n\nFrom Data to Blueprint in Minutes\n2025. 12. | Gemini Agent"
    )

    # 2. Hook
    create_content_slide(
        prs,
        "The Question",
        [
            "매주 1,000개의 신상품이 쏟아지는 시대,",
            "아직도 디자이너의 '감(Feeling)'으로만 싸우시겠습니까?",
            "",
            "패션 산업의 가장 큰 리스크 = '예측 실패로 인한 재고'",
            "소비자는 어제 본 틱톡 트렌드를 오늘 원하지만,",
            "우리의 기획은 3개월 전 데이터에 멈춰 있습니다."
        ]
    )

    # 3. Executive Summary
    create_content_slide(
        prs,
        "Executive Summary",
        [
            "본 제안은 파편화된 정보를 실시간 수집·분석하고,",
            "디자인 시각화부터 생산용 작업지시서까지 자동 생성하는",
            "통합 AI 시스템 구축을 목표로 합니다.",
            "",
            "1. Hyper-Local Trend Sensing: 15+개 채널 실시간 감지",
            "2. AI Merchandiser: Gemini & GLM 기반 멀티모달 분석",
            "3. Production-Ready: 스케치 → 패턴 도면 One-Stop 생성",
            "",
            "기대효과: 기획 리드타임 90% 단축"
        ]
    )

    # 4. Market Context
    create_content_slide(
        prs,
        "Market Context: The Speed War",
        [
            "1. 트렌드 주기의 붕괴 (Micro-Trends)",
            "   - 전통적 시즌(SS/FW) 경계 소멸",
            "   - SNS(Reels, TikTok)발 '초단기 유행' 급부상",
            "",
            "2. 정보의 홍수와 기획 병목",
            "   - Problem: 디자이너 업무 40%가 단순 검색/수집",
            "   - Risk: 데이터 없는 '감' 의존 기획 → 악성 재고"
        ],
        highlight_box="Pain Point:\n\n정보는 넘치는데,\n정작 쓸 만한\n인사이트는 없다."
    )

    # 5. Our Solution Concept
    create_content_slide(
        prs,
        "Our Solution: AI Co-Designer",
        [
            "\"24시간 잠들지 않는 AI 머천다이저를 채용하십시오.\"",
            "",
            "단순 이미지 생성 툴이 아닙니다.",
            "시장 조사부터 기획, 초기 생산 문서까지 대행하는 에이전트입니다.",
            "",
            "Role 1: Trend Watcher (지치지 않는 시장 조사)",
            "Role 2: Data Analyst (소비자 반응 정량 분석)",
            "Role 3: Pattern Maker (생산 고려한 기술 도면 작성)"
        ]
    )

    # 6. Comparison Table (Table Slide)
    create_comparison_table_slide(prs)

    # 7. Process 1 - Sense
    create_content_slide(
        prs,
        "Step 1: Sense (초광각 데이터 수집)",
        [
            "국내 시장 특화 15종 이상 크롤러 가동",
            "",
            "Commerce: 무신사, 29CM (랭킹, 리뷰)",
            "Social: Instagram(해시태그), YouTube(룩북 댓글)",
            "Community (Real Voice):",
            "   - 남성: 펨코, 디매, 이토랜드",
            "   - 여성: 더쿠, 인스티즈, 파우더룸",
            "   - 포털: 네이버 카페/블로그",
            "",
            "차별점: 텍스트 + 이미지 + 댓글 맥락(Context) 동시 수집"
        ]
    )

    # 8. Process 2 - Think
    create_content_slide(
        prs,
        "Step 2: Think (트렌드 인텔리전스)",
        [
            "비정형 데이터를 LLM이 분석하여 '팔리는 키워드' 도출",
            "",
            "Trend Keywords: 급상승 패션 용어 (예: 발레코어)",
            "Sentiment Analysis: 스타일별 대중 호불호 판별",
            "SWOT Report: 아이템별 강점/약점 분석",
            "",
            "Engine: Google Gemini 2.5 (멀티모달) + Zhipu GLM (뉘앙스)"
        ]
    )

    # 9. Process 3 - Create
    create_content_slide(
        prs,
        "Step 3: Create (디자인 시각화)",
        [
            "분석된 데이터 기반 고품질 디자인 시안 생성",
            "",
            "Prompt Engineering: 트렌드 키워드 → 최적 프롬프트 변환",
            "Multi-Style Generation:",
            "   - Lookbook: 실제 모델 착장 실사",
            "   - Flat Lay: 쇼핑몰 상세페이지용 바닥 컷",
            "   - Illustration: 디자이너 감성 스케치"
        ]
    )

    # 10. Process 4 - Realize (Blueprint)
    create_content_slide(
        prs,
        "Step 4: Realize (Blueprint & Tech Pack)",
        [
            "\"그림에서 옷으로\" - 가장 강력한 차별화 기능",
            "",
            "Technical Drawing: 전/후면 도식화 자동 생성",
            "Pattern Pieces: 패턴 조각 분해 및 시각화",
            "Size Spec: KS/ASTM/ISO 표준 기반 치수표 산출",
            "Instructions: 봉제 주의사항 및 원단 요척 계산"
        ],
        highlight_box="Core Value:\n\n단순 이미지가 아닌,\n'작업지시서'를\n만듭니다."
    )

    # 11. Architecture
    create_content_slide(
        prs,
        "System Architecture",
        [
            "확장성과 보안을 고려한 모던 아키텍처",
            "",
            "Frontend: 직관적 웹 대시보드 (HTML5/JS)",
            "Backend: FastAPI (Python) - 고성능 비동기 처리",
            "AI Core: LangChain Orchestrator (Gemini/GLM/Flux)",
            "Data Pipeline: Playwright/Selenium → Vector DB"
        ]
    )

    # 12. Roadmap
    create_content_slide(
        prs,
        "Implementation Roadmap (3개월)",
        [
            "Phase 1: Foundation (Month 1)",
            "   - 핵심 크롤러 5종 및 데이터 파이프라인 구축",
            "   - Goal: \"데이터가 흐르기 시작한다\"",
            "",
            "Phase 2: Intelligence & Generation (Month 2)",
            "   - 프롬프트 최적화 및 디자인 생성 고도화",
            "   - Goal: \"원하는 디자인이 나온다\"",
            "",
            "Phase 3: Production Ready (Month 3)",
            "   - 블루프린트(도면) 알고리즘 정밀화, UI 폴리싱",
            "   - Goal: \"실무에 바로 쓴다\""
        ]
    )

    # 13. Expected Impact
    create_content_slide(
        prs,
        "Expected Impact",
        [
            "1. Speed to Market",
            "   - 기획~샘플 발주: 2주 → 2일 단축",
            "",
            "2. Hit Ratio",
            "   - 데이터 기반 디자인으로 베스트셀러 적중률 향상",
            "",
            "3. Cost Saving",
            "   - 불필요한 샘플 제작 감소 = 비용 절감",
            "",
            "4. Asset Management",
            "   - 사내 트렌드 DB 및 디자인 자산 축적"
        ]
    )

    # 14. Conclusion
    create_title_slide(
        prs,
        "Fashion is Data.",
        "지금 바로, 패션 R&D의 혁신을 시작하십시오.\n\n감사합니다."
    )

    output_path = "Fashion_AI_System_Proposal_v2.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    create_presentation()
