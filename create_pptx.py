from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Define content
    slides_content = [
        {
            "layout": "title",
            "title": "AI 기반 패션 트렌드 분석 및\n디자인 자동 생성 시스템 구축 제안",
            "subtitle": "Data-Driven Fashion Design & Trend Forecasting System\n\n2025. 12.\n제안자: Gemini Agent"
        },
        {
            "layout": "content",
            "title": "요약 (Executive Summary)",
            "body": "데이터에서 디자인까지, 패션 R&D의 혁신\n\n1. 배경\n   - 패션 트렌드 주기의 단축 (Ultra-Fast Fashion)\n   - 데이터 과잉으로 인한 의사결정 지연\n\n2. 핵심 솔루션\n   - 초광각 데이터 수집: SNS, 커뮤니티, 뉴스 등\n   - AI 트렌드 인텔리전스: LLM 기반 감성 분석 및 키워드 추출\n   - One-Stop 디자인 생성: 트렌드 반영 디자인 → 도면(Blueprint) 자동화\n\n3. 기대 효과\n   - 기획 단계 소요 시간 70% 단축\n   - 데이터 기반 의사결정으로 적중률 향상"
        },
        {
            "layout": "content",
            "title": "시장 환경 분석 및 Pain Points",
            "body": "1. 급변하는 패션 트렌드\n   - 기존 시즌제(SS/FW) 기획 방식의 한계\n   - SNS 숏폼 콘텐츠에 의한 '마이크로 트렌드' 급부상\n\n2. 데이터의 홍수와 분석의 어려움\n   - Problem: 참고해야 할 사이트와 정보량 폭증\n   - Pain Point: 디자이너 업무의 40%가 단순 정보 수집\n\n3. 디자인 프로세스의 병목\n   - 아이디어 스케치에서 샘플링까지 긴 리드타임\n   - 단순 변형(Variation) 작업을 위한 반복 업무"
        },
        {
            "layout": "content",
            "title": "To-Be: AI 기반 통합 디자인 시스템",
            "body": "Trend Watching to Product Creation\n\n[System Concept]\n1. Sense (감지): 15+개 채널 실시간 모니터링 (크롤러 군단)\n2. Think (분석): 소비자 반응 분석 및 인사이트 도출 (LLM Agent)\n3. Create (생성): 분석된 키워드로 디자인 및 도면 생성 (Image Gen AI)\n\n[Core Value]\n- Speed: 트렌드 포착 후 디자인 시안까지 10분 이내\n- Accuracy: 정량적 데이터 기반의 디자인 제안\n- Extension: 생산 연계형 도면(Blueprint) 제공"
        },
        {
            "layout": "content",
            "title": "System Architecture",
            "body": "안정적이고 확장 가능한 모듈형 아키텍처\n\n1. User Interface\n   - 웹 대시보드, 분석 리포트 뷰어, 디자인 에디터\n\n2. Service Layer\n   - Trend Analyzer (키워드/리포트), Design Generator (이미지)\n   - Blueprint Maker (도면/치수)\n\n3. AI Layer\n   - LLM: Google Gemini 2.5, Zhipu GLM-4\n   - Vision: Stable Diffusion, Nano Banana\n\n4. Data Layer\n   - Crawlers (SNS/News), Vector DB"
        },
        {
            "layout": "content",
            "title": "핵심 기능 1 - 멀티 채널 트렌드 수집",
            "body": "Deep-Dive: 광범위 데이터 수집 (The Eyes)\n\n1. 소셜 미디어\n   - Instagram: 해시태그 기반 스타일/반응 수집\n   - YouTube: 패션 룩북 영상 댓글 분석\n\n2. 커뮤니티 (버티컬)\n   - 무신사/29CM: 실시간 랭킹, 리뷰 데이터\n   - 패션 커뮤니티: 고인물/더쿠/펨코 등 여론 추적\n\n3. 차별점\n   - 텍스트뿐만 아니라 '맥락'과 '이미지'를 함께 수집하여 멀티모달 분석"
        },
        {
            "layout": "content",
            "title": "핵심 기능 2 - AI 트렌드 인텔리전스",
            "body": "Deep-Dive: 트렌드 분석 및 인사이트 (The Brain)\n\n1. 키워드 추출\n   - 급상승 패션 키워드 (예: 고프코어, 발레코어) 도출\n\n2. 감성 분석 (Sentiment Analysis)\n   - 특정 스타일/브랜드에 대한 긍정/부정 여론 판별\n\n3. 트렌드 리포트 생성\n   - LLM이 '2025 F/W 아우터 트렌드' 등 주제별 리포트 자동 작성"
        },
        {
            "layout": "content",
            "title": "핵심 기능 3 - 디자인 및 도면 자동화",
            "body": "Deep-Dive: 디자인 & 도면 생성 (The Hands)\n\n1. 프롬프트 자동 최적화\n   - 트렌드 키워드를 고품질 이미지 프롬프트로 변환\n\n2. 다양한 스타일 생성\n   - Realism (실사), Illustration (스케치), Flat Lay (바닥샷)\n\n3. 블루프린트(Blueprint) 생성\n   - 디자인 기반 기술 도면(도식화) 자동 생성\n   - 제작 지시 사항 및 예상 치수표(Size Spec) 제안"
        },
        {
            "layout": "content",
            "title": "개발 일정 (3개월 로드맵)",
            "body": "Phase 1: 기반 구축 및 커스터마이징 (Month 1)\n- 요구사항 분석, 크롤링 타겟 확정\n- 데이터 파이프라인 최적화 (Milestone: 대시보드 오픈)\n\nPhase 2: AI 모델 고도화 및 로직 구현 (Month 2)\n- 프롬프트 엔지니어링 및 파인튜닝\n- 도면 생성 알고리즘 정교화 (Milestone: 디자인/도면 테스트)\n\nPhase 3: UI/UX 통합 및 안정화 (Month 3)\n- 웹 인터페이스 고도화\n- 통합 테스트 및 인수인계 (Milestone: 최종 오픈)"
        },
        {
            "layout": "content",
            "title": "기대 효과 및 결론",
            "body": "Expected Benefits\n1. 기획 효율성 300% 증대 (1주일 → 1일)\n2. 적중률 높은 디자인 (데이터 기반 의사결정)\n3. 자산화 (사내 트렌드 DB 및 디자인 아카이브 축적)\n\nConclusion\n단순한 툴이 아닌, 'AI 보조 디자이너'를 채용하여\n크리에이티브 디렉터가 본질적인 창작에 집중할 수 있는 환경 제안"
        }
    ]

    for slide_info in slides_content:
        if slide_info["layout"] == "title":
            slide_layout = prs.slide_layouts[0] # Title Slide
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = slide_info["title"]
            subtitle.text = slide_info["subtitle"]
            
        elif slide_info["layout"] == "content":
            slide_layout = prs.slide_layouts[1] # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            body = slide.placeholders[1]
            
            title.text = slide_info["title"]
            body.text = slide_info["body"]

            # Basic formatting for body text
            for paragraph in body.text_frame.paragraphs:
                paragraph.font.size = Pt(18)
                paragraph.space_after = Pt(10)

    output_path = "Fashion_AI_System_Proposal.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    create_presentation()
